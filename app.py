import streamlit as st
import fitz  # PyMuPDF
import io
from pptx import Presentation

st.title("PDF to PPT Converter")
st.write("Upload a PDF file to convert each page into a slide in a PowerPoint presentation.")

# File uploader accepts PDF files
pdf_file = st.file_uploader("Choose a PDF file", type=["pdf"])

if pdf_file is not None:
    with st.spinner("Converting PDF to PPT..."):
        # Read the uploaded PDF file into memory
        pdf_bytes = pdf_file.read()
        
        # Open the PDF with PyMuPDF
        doc = fitz.open("pdf", pdf_bytes)
        
        # Create a new PowerPoint presentation
        prs = Presentation()
        
        # Process each page of the PDF
        for page in doc:
            # Render the page as an image at 300 dpi for quality
            pix = page.get_pixmap(dpi=300)
            img_bytes = pix.tobytes("png")
            
            # Create an in-memory image stream
            image_stream = io.BytesIO(img_bytes)
            
            # Add a blank slide (layout index 6 is typically blank)
            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # Set the image to fill the entire slide
            slide_width = prs.slide_width
            slide_height = prs.slide_height
            slide.shapes.add_picture(image_stream, 0, 0, width=slide_width, height=slide_height)
        
        # Save the presentation to an in-memory bytes buffer
        ppt_io = io.BytesIO()
        prs.save(ppt_io)
        ppt_io.seek(0)
    
    st.success("Conversion complete!")
    st.download_button(
        label="Download PPTX",
        data=ppt_io,
        file_name="output.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
